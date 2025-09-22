# core/controller.py

import tkinter as tk
from tkinter import messagebox
import os
from pynput import keyboard
import docx
import fitz
from PIL import Image
import threading
from pptx import Presentation
import json
import sys

# 导入我们自己的模块
from features.floating_ball import FloatingBall
from features.instructions_window import InstructionsWindow
from features.settings_window import SettingsWindow
from features.tray_icon import TrayIcon
from .config_manager import ConfigManager
from .ui import ScreenshotTaker, ResultWindow
from .utils import log, set_proxy, configure_gemini, get_screen_scaling_factor

# 主控制器 
class MainController:
    def __init__(self, root):
        log("[LOG-C1] MainController 初始化。")
        self.listener = None 
        self.config_manager = ConfigManager()
        self.config = None
        self.gemini_model = None
        self.scaling_factor = 1.0
        self.current_pressed = set()
        self.hotkey_actions = {} 
        self.is_running_action = False
        self.root = root
        self.floating_ball = None
        self.tray_icon = None

    def open_instructions_window(self):
        """打开使用说明窗口"""
        if self.floating_ball:
            self.floating_ball.reset_idle_timer()
        InstructionsWindow(self.root)

    def restart_app(self):
        """重启应用程序"""
        log("正在准备重启应用程序...")
        
        try:
            # 在退出当前进程之前，先启动一个新进程
            # sys.executable 指的是当前运行的Python解释器
            # sys.argv[0] 指的是我们程序的主入口文件 (main.py)
            
            # 为了确保新进程能找到所有模块，我们最好设置工作目录
            # app_dir 的计算逻辑与 ConfigManager 类似
            if getattr(sys, 'frozen', False):
                app_dir = os.path.dirname(sys.executable)
                main_script = sys.executable
            else:
                app_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
                main_script = os.path.join(app_dir, 'main.py')

            log(f"启动新进程: {sys.executable} {main_script}")
            
            # 使用 subprocess.Popen 启动一个完全分离的新进程
            import subprocess
            subprocess.Popen([sys.executable, main_script], cwd=app_dir)

            # 启动新进程后，立即干净地退出当前进程
            log("新进程已启动，正在退出当前进程...")
            self.exit_app()

        except Exception as e:
            log(f"重启失败: {e}")
            messagebox.showerror("重启失败", f"无法重启应用程序:\n{e}")

    def _start_task(self, task_name="新任务"):
        """
        所有任务的统一入口。负责检查状态、设置代理和更新标志。
        如果可以开始任务，返回 True；如果当前正忙，返回 False。
        """
        if self.is_running_action:
            log(f"警告：当前已有任务在运行，任务 '{task_name}' 的触发被忽略。")
            messagebox.showwarning("忙碌中", "请等待上一个任务完成。")
            return False
        
        log(f"---------- {task_name} 开始 ----------")
        # 1. 重新检查并设置代理
        set_proxy(self.config.get("proxy_url", ""))
        
        # 2. 设置运行状态标志
        self.is_running_action = True
        return True

    def _end_task(self):
        """所有任务的统一结束点。负责重置状态标志。"""
        self.is_running_action = False
        log("任务结束，返回待机状态。")
    def on_press(self, key):
        # 如果有任务在运行，直接忽略任何按键，防止冲突
        if self.is_running_action:
            return
            
        try:
            normalized_key = keyboard.KeyCode.from_char(key.char.lower())
        except AttributeError:
            normalized_key = key
        
        self.current_pressed.add(normalized_key)

        # 遍历所有已注册的动作快捷键
        for action_name, key_set in self.hotkey_actions.items():
            if self.current_pressed == key_set:
                log(f"动作 '{action_name}' 的快捷键被触发。")
                if self.root:
                    self.root.after(0, self.trigger_action, action_name)
                break
            
    def on_release(self, key):
        """按键释放的监听回调，只负责清理 (最终绝对稳定版)"""
        # 释放任何一个键都重置当前按键组合状态
        self.current_pressed.clear()
        #log(f"按键释放，重置快捷键状态。")

    def setup_from_config(self):
        """根据加载的配置初始化应用，失败则返回False"""
        # 提前设置代理，并获取返回的代理地址
        proxy_url = set_proxy(self.config.get("proxy_url", ""))

        model_name = self.config.get("model_name", "gemini-1.5-flash-latest")
        
        # 将获取到的代理地址传给 configure_gemini
        self.gemini_model = configure_gemini(
            api_key=self.config.get("api_key", ""), 
            model_name=model_name,
            proxy_url=proxy_url
        )

        if not self.gemini_model:
            self.show_error_and_exit(f"Gemini配置错误: API Key 无效、网络问题或模型名称 '{model_name}' 不正确。\n\n请检查 config.json 文件和代理设置。\n如果更新了API key，请重启程序以生效。")
            return False

        # 先把 actions 从 config 中取出来
        actions = self.config.get("actions", {})
        if not actions:
            self.show_error_and_exit("配置文件中未找到 'actions' 配置项。")
            return False
            
        log("正在解析 actions...")
        for action_name, details in actions.items():
            hotkey_str = details.get("hotkey")
            if not hotkey_str:
                log(f"警告: 动作 '{action_name}' 没有配置 hotkey，将被忽略。")
                continue
            
            key_set = self.parse_hotkey(hotkey_str)
            if not key_set:
                self.show_error_and_exit(f"快捷键配置错误: 无法解析动作 '{action_name}' 的快捷键 '{hotkey_str}'。")
                return False
            
            self.hotkey_actions[action_name] = key_set
            log(f"  -> 已注册动作: '{action_name}', 快捷键: {hotkey_str}")

        if not self.hotkey_actions:
            self.show_error_and_exit("未成功注册任何有效的快捷键动作。")
            return False
            
        self.scaling_factor = get_screen_scaling_factor()
        return True

    def parse_hotkey(self, hotkey_str):
        """将字符串 'a+b+c' 解析为 pynput 按键集合 (最终修正版)"""
        if not isinstance(hotkey_str, str) or not hotkey_str: return set()
        
        key_map = {
            'ctrl': keyboard.Key.ctrl, 'alt': keyboard.Key.alt,
            'shift': keyboard.Key.shift, 'cmd': keyboard.Key.cmd, 'win': keyboard.Key.cmd
        }
        keys = set()
        
        parts = hotkey_str.lower().split('+')
        for part in parts:
            part = part.strip()
            if part in key_map:
                keys.add(key_map[part])
            elif len(part) == 1:
                keys.add(keyboard.KeyCode.from_char(part))
            else:
                log(f"无法识别的快捷键部分: '{part}'")
                return set() # 返回空集合表示解析失败
        
        return keys

    def trigger_action(self, action_name):
        """根据快捷键动作名称触发相应的流程"""
        if not self._start_task(f"快捷键任务 ({action_name})"):
            return # 如果任务启动失败（正忙），则直接返回

        # 如果任务成功启动，则继续执行
        if action_name == "screenshot":
            self.start_screenshot_flow()
        elif action_name == "clipboard_text":
            self.start_clipboard_flow()
        else:
            log(f"错误: 未知的动作名称 '{action_name}'")
            self._end_task() # 未知动作，立即结束任务

    def start_clipboard_flow(self):
        """处理剪贴板文本的流程"""
        try:
            clipboard_content = self.root.clipboard_get()
            if not clipboard_content.strip():
                messagebox.showinfo("提示", "剪贴板内容为空。")
                self._end_task() # 流程结束，重置标志
                return
            
            action_config = self.config["actions"]["clipboard_text"]
            prompt = action_config.get("prompt", "请处理这段文本:") # 默认值
            
            log("成功获取剪贴板文本，准备显示结果窗口。")
            self.show_result_window(
                prompt=prompt, 
                task_type="text", 
                task_data=clipboard_content
            )

        except tk.TclError:
            messagebox.showinfo("提示", "无法从剪贴板获取文本内容。")
            self._end_task() # 流程结束，重置标志
        except Exception as e:
            messagebox.showerror("错误", f"处理剪贴板时发生未知错误: {e}")
            log(f"处理剪贴板时发生未知错误: {e}")
            self._end_task() # 流程结束，重置标志

    def open_settings_window(self):
        """打开设置窗口"""
        if self.floating_ball:
            self.floating_ball.reset_idle_timer()

        # 创建设置窗口实例，传入当前配置和保存回调
        SettingsWindow(self.root, self.config, self.save_config_and_update)

    def save_config_and_update(self, new_config):
        """
        接收来自设置窗口的新配置，更新内存并写入文件。
        """
        # 1. 更新内存中的配置
        self.config = new_config
        log("内存中的配置已更新。")

        # 2. 持久化到 config.json 文件
        try:
            # config_manager.filepath 包含了正确的、带路径的文件名
            filepath = self.config_manager.filepath
            with open(filepath, 'w', encoding='utf-8') as f:
                json.dump(new_config, f, indent=2, ensure_ascii=False)
            log(f"配置已成功写入到 {filepath}")
        except Exception as e:
            log(f"错误：无法写入配置文件: {e}")
            messagebox.showerror("保存失败", f"无法将设置写入到 config.json 文件:\n{e}")
    
    def create_widgets(self):
        """创建所有UI组件"""
        
        # --- 创建包装后的回调函数 ---
        def start_chat_and_close_menu():
            if self.floating_ball: self.floating_ball.close_menu()
            self.start_temporary_chat()

        def hide_ball_and_close_menu():
            if self.floating_ball: self.floating_ball.close_menu()
            self.hide_ball_to_tray()

        def open_settings_and_close_menu():
            if self.floating_ball: self.floating_ball.close_menu()
            self.open_settings_window()
            
        def open_instructions_and_close_menu(): # <--- 新增这个包装函数
            if self.floating_ball: self.floating_ball.close_menu()
            self.open_instructions_window()

        def restart_and_close_menu():
            if self.floating_ball: self.floating_ball.close_menu()
            self.restart_app()

        def exit_and_close_menu():
            if self.floating_ball: self.floating_ball.close_menu()
            self.exit_app()

        # --- 实例化悬浮球，并传入所有回调 ---
        self.floating_ball = FloatingBall(
            master=self.root,
            on_start_chat_callback=start_chat_and_close_menu,
            on_hide_callback=hide_ball_and_close_menu,
            on_drop_callback=self.handle_drop_data,
            on_settings_callback=open_settings_and_close_menu,
            on_instructions_callback=open_instructions_and_close_menu, # <--- 新增这一行
            on_restart_callback=restart_and_close_menu,
            on_exit_callback=exit_and_close_menu
        )
        
        # --- 实例化托盘图标 (保持不变) ---
        self.tray_icon = TrayIcon(
            on_show_callback=self.show_ball_from_tray,
            on_exit_callback=self.exit_app
        )


    def start_temporary_chat(self):
        """由悬浮球触发，启动一个临时的纯文本聊天"""
        if self.floating_ball:
            self.floating_ball.reset_idle_timer()

        if not self._start_task("悬浮球新对话"):
            return

        prompt = "你好！"
        self.show_result_window(
            prompt=prompt,
            task_type="text",
            task_data=""
        )

    def hide_ball_to_tray(self):
        """隐藏悬浮球并显示托盘图标"""
        if self.floating_ball:
            self.floating_ball.reset_idle_timer()
        if self.floating_ball:
            self.floating_ball.hide()
        if self.tray_icon:
            self.tray_icon.start() # 仅在需要时启动托盘线程

    def show_ball_from_tray(self):
        """从托盘恢复悬浮球"""
        if self.tray_icon:
            self.tray_icon.stop()
            # pystray的stop是非阻塞的，需要一点时间来确保线程退出
            self.root.after(100, self._show_ball_action)
            
    def _show_ball_action(self):
        """确保在托盘退出后显示悬浮球"""
        if self.floating_ball:
            self.floating_ball.show()
            self.floating_ball.window.lift() # 确保窗口在最顶层
            self.floating_ball.window.focus_force()

    def exit_app(self):
        """干净地退出整个应用程序"""
        log("正在退出应用程序...")
        
        # --- 新增的关键代码 ---
        if self.listener and self.listener.is_alive():
            log("正在停止键盘监听器...")
            self.listener.stop()
            # (可选) self.listener.join() # 可以等待线程完全停止，但通常stop()就够了
        # --- 新增代码结束 ---
        
        if self.tray_icon:
            self.tray_icon.stop()
        if self.floating_ball:
            self.floating_ball.destroy()
        if self.root:
            self.root.quit()
            # self.root.destroy() # quit() 之后 mainloop 结束，destroy() 可能会引发错误
        
        log("应用程序已退出。")
        # os._exit(0) 应该只作为最后的手段，现在我们有了更优雅的退出方式，可以先注释掉它
        os._exit(0)

    # 我们需要修改 show_error_and_exit
    def show_error_and_exit(self, message):
        """显示一个错误消息框，然后退出程序"""
        log(f"启动错误: {message.replace(os.linesep, ' ')}")
        # 错误可能在主循环开始前发生，此时 messagebox 可能无法正常显示
        # 创建一个临时root来显示错误
        temp_root = tk.Tk()
        temp_root.withdraw()
        messagebox.showerror("启动错误", message, parent=temp_root)
        temp_root.destroy()
        self.exit_app()
    def run(self):
        """启动主程序"""
        log("[LOG-C3] MainController.run() 启动。")
        # self.root = tk.Tk() <--- 删除
        # self.root.withdraw() <--- 删除
        
        self.config = self.config_manager.load_config()

        if not self.config or not self.config.get("api_key"):
            self.show_error_and_exit(
                "配置文件 config.json 不存在或无效！\n\n"
                "请在程序同级目录下创建 config.json 文件，\n"
                "并填入您的 API Key 和其他配置。\n"
                "API Key 可以在 https://aistudio.google.com/app/apikey 免费获取。\n"
                "如果更新了API key，请重启程序以生效。"
            )
            return

        if not self.setup_from_config():
            return

        log("正在创建悬浮球和系统托盘...")
        self.create_widgets()

        log("Gemini助手已启动，在后台等待快捷键...")
        log("已注册的快捷键动作如下:")
        actions_config = self.config.get("actions", {})
        for name in self.hotkey_actions.keys():
            hotkey_str = actions_config.get(name, {}).get("hotkey", "未知")
            log(f"  -> {name}: {hotkey_str.upper()}")
        

        self.listener = keyboard.Listener(on_press=self.on_press, on_release=self.on_release)
        self.listener.start()
        
        # 我们需要在 run 的最后确保 WM_DELETE_WINDOW 协议被设置
        self.root.protocol("WM_DELETE_WINDOW", self.exit_app)
        self.root.mainloop()


    def start_screenshot_flow(self):
        """处理截图的流程"""
        log("[LOG-C6] 准备创建 ScreenshotTaker 实例...")
        screenshot_path = os.path.join(self.config_manager.app_dir, "temp_screenshot.png")
        # 将截图路径作为额外信息传递给回调函数
        callback = lambda cancelled=False: self.screenshot_done(cancelled, screenshot_path)
        ScreenshotTaker(self.scaling_factor, callback, screenshot_path)
        log("[LOG-C7] ScreenshotTaker 实例创建完成。")
    
    def screenshot_done(self, cancelled=False, screenshot_path=None):
        if cancelled:
            log("截图任务被取消。")
            self._end_task() # 使用新的结束方法
        elif screenshot_path:
            action_config = self.config["actions"]["screenshot"]
            prompt = action_config.get("prompt", "请描述这张图片:")
            self.show_result_window(
                prompt=prompt,
                task_type="image",
                task_data=screenshot_path
            )
        else:
            log("截图完成，但未提供截图路径，任务异常结束。")
            self._end_task() # 异常情况也要结束任务
        
    def show_result_window(self, prompt, task_type, task_data):
        log("[LOG-C8] show_result_window 被调用。")
        result_win = ResultWindow(
            model=self.gemini_model, 
            config_manager=self.config_manager,
            prompt=prompt,
            task_type=task_type,
            task_data=task_data
        )
        result_win.protocol("WM_DELETE_WINDOW", lambda: self.on_result_window_close(result_win))
        
    def on_result_window_close(self, window):
        window.destroy()
        self._end_task()

    def handle_drop_data(self, data):
        """拖放数据的总入口和调度中心。"""
        if not self._start_task("悬浮球拖放任务"):
            return

        cleaned_data = data.strip()
        if cleaned_data.startswith('{') and cleaned_data.endswith('}'):
            cleaned_data = cleaned_data[1:-1]

        if os.path.exists(cleaned_data):
            self.process_dropped_files([cleaned_data])
        else:
            self.process_dropped_text(data)

    def process_dropped_text(self, text_content):
        """处理拖放的纯文本"""
        log("处理拖放的纯文本...")
        # self.is_running_action = True <--- 删除这行，_start_task 已经处理
        
        action_config = self.config["actions"].get("clipboard_text", {})
        prompt = action_config.get("prompt", "请处理以下文本:")
        
        self.show_result_window(
            prompt=prompt,
            task_type="text",
            task_data=text_content
        )

    def process_dropped_files(self, filepaths):
        """处理拖放的文件列表，并根据文件类型分发任务。"""
        if not filepaths:
            self._end_task() # 如果没有文件路径，直接结束任务
            return
            
        filepath = filepaths[0]
        log(f"处理拖放的文件: {filepath}")
        # self.is_running_action = True <--- 删除这行，_start_task 已经处理

        ext = os.path.splitext(filepath)[1].lower()
        
        # --- 智能选择 Prompt 和处理方式 ---
        
        drop_handlers_config = self.config.get("drop_handlers", {})
        if ext in drop_handlers_config:
            handler_info = drop_handlers_config[ext]
            prompt = handler_info.get("prompt", f"请处理这个 {ext} 文件。")
            log(f"使用为 '{ext}' 类型定义的专属 prompt。")
            
            if ext == ".docx":
                self.handle_docx_file(filepath, prompt)
            elif ext == ".pdf":
                self.handle_pdf_file(filepath, prompt)
            elif ext == ".pptx": 
                self.handle_pptx_file(filepath, prompt)
            elif ext in [".py", ".js", ".html", ".css", ".java", ".cpp"]:
                self.handle_text_based_file(filepath, prompt)
            else:
                self.handle_text_based_file(filepath, prompt)
            return

        if ext in ['.png', '.jpg', '.jpeg', '.bmp', '.webp']:
            log("回退到通用的图片处理逻辑。")
            self.handle_dropped_image(filepath)
        elif ext in ['.txt', '.md', '.json']:
            log("回退到通用的文本文件处理逻辑。")
            self.handle_text_based_file(filepath, None)
        else:
            unsupported_message = f"不支持的文件类型: {ext}\n\n当前支持图片、文本和部分文档格式。"
            log(unsupported_message)
            self._end_task()
            messagebox.showinfo("操作失败", unsupported_message)

    def handle_pdf_file(self, filepath, prompt):
        """
        具体处理 .pdf 文件的逻辑。
        采用图文混合模式，提取全部文本和页面截图。
        """
        log(f"正在处理PDF文件: {filepath}")
        
        # --- 创建一个后台线程来处理耗时的PDF解析 ---
        def pdf_processing_thread():
            try:
                doc = fitz.open(filepath)
                
                full_text = []
                page_images = []
                
                log(f"PDF共有 {len(doc)} 页。开始提取内容...")
                
                # 为了防止内容过长和处理时间过久，可以设置一个最大处理页数
                MAX_PAGES_TO_PROCESS = 50 
                pages_to_process = doc.page_count
                if pages_to_process > MAX_PAGES_TO_PROCESS:
                    log(f"警告：PDF页数过多，仅处理前 {MAX_PAGES_TO_PROCESS} 页。")
                    pages_to_process = MAX_PAGES_TO_PROCESS

                for i, page in enumerate(doc.load_page(page_num) for page_num in range(pages_to_process)):
                    # 1. 提取文本
                    full_text.append(page.get_text("text"))
                    
                    # 2. 渲染页面为图片
                    # 设置DPI以提高图片质量，150是个不错的折中值
                    pix = page.get_pixmap(dpi=150)
                    # 从pixmap的样本数据创建Pillow Image对象
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    page_images.append(img)
                    
                    log(f"  ...已处理第 {i+1}/{pages_to_process} 页")

                doc.close()

                # 将所有内容打包
                content = "\n".join(full_text)
                
                # 将任务调度回主线程以显示结果窗口
                # 注意：这里我们将图片列表和文本一起传递
                self.root.after(0, self.show_result_window_for_multimodal, 
                                prompt, "pdf_multimodal", (content, page_images))

            except Exception as e:
                log(f"解析 .pdf 文件失败: {filepath}, 错误: {e}")
                # 确保在主线程中显示错误并重置状态
                def show_error_and_reset():
                    messagebox.showerror("PDF解析失败", f"无法解析 .pdf 文件:\n{e}")
                    self._end_task()
                self.root.after(0, show_error_and_reset)

        # 启动后台处理线程
        threading.Thread(target=pdf_processing_thread, daemon=True).start()

    def handle_text_based_file(self, filepath, prompt=None):
        """
        统一处理所有基于文本的文件（.txt, .py, .md 等）。
        如果 prompt 为 None，则复用剪贴板的 prompt。
        """
        if prompt is None:
            # 复用旧逻辑
            action_config = self.config["actions"].get("clipboard_text", {})
            prompt = action_config.get("prompt", "请分析以下文件内容:")
            
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                content = f.read()
        except Exception as e:
            log(f"读取文件失败: {filepath}, 错误: {e}")
            messagebox.showerror("读取失败", f"无法读取文件内容:\n{e}")
            self._end_task()
            return
            
        self.show_result_window(
            prompt=prompt,
            task_type="text",
            task_data=content
        )

    def handle_docx_file(self, filepath, prompt):
        """具体处理 .docx 文件的逻辑"""
        log(f"正在从 .docx 文件提取文本: {filepath}")
        try:
            doc = docx.Document(filepath)
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            
            # (可选) 提取表格内容
            # for table in doc.tables:
            #     for row in table.rows:
            #         for cell in row.cells:
            #             full_text.append(cell.text)

            content = '\n'.join(full_text)

        except Exception as e:
            log(f"解析 .docx 文件失败: {filepath}, 错误: {e}")
            messagebox.showerror("解析失败", f"无法解析 .docx 文件:\n{e}")
            self._end_task()
            return

        self.show_result_window(
            prompt=prompt,
            task_type="text", # docx内容也是文本
            task_data=content
        )

    def handle_pptx_file(self, filepath, prompt):
        """
        具体处理 .pptx 文件的逻辑。
        提取所有幻灯片的文本内容。
        """
        log(f"正在从 .pptx 文件提取文本: {filepath}")
        try:
            prs = Presentation(filepath)
            full_text = []
            
            # 设置一个最大处理幻灯片数量的上限
            MAX_SLIDES_TO_PROCESS = 100
            slides_to_process = len(prs.slides)
            if slides_to_process > MAX_SLIDES_TO_PROCESS:
                log(f"警告：PPT页数过多，仅处理前 {MAX_SLIDES_TO_PROCESS} 页。")
                slides_to_process = MAX_SLIDES_TO_PROCESS
            
            log(f"PPT共有 {len(prs.slides)} 页。开始提取内容...")

            for i, slide in enumerate(prs.slides):
                if i >= slides_to_process:
                    break
                
                # 为每一页添加一个标题，让AI更好地理解结构
                full_text.append(f"\n--- 幻灯片 {i + 1} ---\n")
                
                # 遍历幻灯片中的所有形状 (shape)
                for shape in slide.shapes:
                    # 检查形状是否包含文本框
                    if not shape.has_text_frame:
                        continue
                    # 遍历文本框中的所有段落
                    for paragraph in shape.text_frame.paragraphs:
                        # 遍历段落中的所有文本块 (run)
                        for run in paragraph.runs:
                            full_text.append(run.text)
                log(f"  ...已处理第 {i+1}/{slides_to_process} 页")
            
            content = '\n'.join(full_text)

        except Exception as e:
            log(f"解析 .pptx 文件失败: {filepath}, 错误: {e}")
            messagebox.showerror("解析失败", f"无法解析 .pptx 文件:\n{e}")
            self._end_task()
            return

        self.show_result_window(
            prompt=prompt,
            task_type="text", # .pptx 的内容也被处理为纯文本
            task_data=content
        )

    def handle_dropped_image(self, filepath):
        """具体处理图片文件的逻辑"""
        action_config = self.config["actions"].get("screenshot", {})
        prompt = action_config.get("prompt", "请描述这张图片:")
        
        self.show_result_window(
            prompt=prompt,
            task_type="image_from_path",
            task_data=filepath
        )
    def show_result_window_for_multimodal(self, prompt, task_type, task_data_tuple):
        """专门为多模态任务（如PDF）调用结果窗口"""
        log("[LOG-C8-MM] show_result_window_for_multimodal 被调用。")
        result_win = ResultWindow(
            model=self.gemini_model, 
            config_manager=self.config_manager,
            prompt=prompt,
            task_type=task_type,
            task_data=task_data_tuple # task_data现在是一个元组 (text, [image1, image2, ...])
        )
        result_win.protocol("WM_DELETE_WINDOW", lambda: self.on_result_window_close(result_win))